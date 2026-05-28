#!/usr/bin/env python3
"""Build APEX Autopilot investor pitch deck as PowerPoint (.pptx)."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "APEX_Autopilot_Pitch_Deck.pptx"

# Brand palette (Bold Signal–inspired)
INK = RGBColor(247, 248, 251)
MUTED = RGBColor(170, 181, 202)
ACCENT = RGBColor(255, 106, 26)
ACCENT2 = RGBColor(45, 157, 255)
BG = RGBColor(11, 15, 23)
PANEL = RGBColor(16, 23, 37)
OK = RGBColor(47, 230, 167)


def _set_slide_bg(slide, color: RGBColor = BG) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_textbox(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str,
    *,
    size: int = 18,
    bold: bool = False,
    color: RGBColor = INK,
    align: PP_ALIGN = PP_ALIGN.LEFT,
) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Calibri"


def _add_bullets(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    lines: list[str],
    *,
    size: int = 16,
    color: RGBColor = MUTED,
) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(8)


def _accent_bar(slide) -> None:
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE — use integer to avoid extra import
        Inches(0),
        Inches(0),
        Inches(0.12),
        Inches(7.5),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()


def build() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 1 — Title
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s)
    _accent_bar(s)
    _add_textbox(s, 0.7, 0.55, 10, 0.4, "INSTITUTIONAL PAPER TRADING PLATFORM", size=14, bold=True, color=ACCENT)
    _add_textbox(s, 0.7, 1.1, 11, 1.2, "APEX Autopilot", size=54, bold=True)
    _add_textbox(
        s,
        0.7,
        2.35,
        11,
        1.2,
        "Autonomous prediction-market intelligence and execution for Kalshi + Polymarket "
        "with risk gates, operator controls, and full auditability.",
        size=20,
        color=MUTED,
    )
    _add_textbox(s, 0.7, 4.0, 3.5, 0.9, "100+", size=36, bold=True, color=INK)
    _add_textbox(s, 0.7, 4.75, 3.5, 0.5, "Live arb opportunities surfaced", size=14, color=MUTED)
    _add_textbox(s, 4.5, 4.0, 3.5, 0.9, "L0–L4", size=36, bold=True, color=INK)
    _add_textbox(s, 4.5, 4.75, 3.5, 0.5, "Layered architecture in production test", size=14, color=MUTED)
    _add_textbox(s, 8.3, 4.0, 3.5, 0.9, "260 Days", size=36, bold=True, color=INK)
    _add_textbox(s, 8.3, 4.75, 3.5, 0.5, "Execution runbook with milestones", size=14, color=MUTED)

    # 2 — Problem
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s)
    _accent_bar(s)
    _add_textbox(s, 0.7, 0.55, 10, 0.4, "THE PROBLEM", size=14, bold=True, color=ACCENT)
    _add_textbox(s, 0.7, 1.0, 11.5, 1.0, "Human traders miss speed, context, and discipline.", size=36, bold=True)
    _add_bullets(
        s,
        0.7,
        2.3,
        11.5,
        4.5,
        [
            "Fragmented data — Kalshi, Polymarket, macro feeds, and settlement sources are slow to reconcile.",
            "Execution lag — arbitrage windows collapse in seconds while manual workflows still evaluate assumptions.",
            "Inconsistent risk — without deterministic gating, teams drift from policy and forensics weaken.",
            "Result: alpha decay, avoidable risk, and poor reproducibility.",
        ],
        size=17,
    )

    # 3 — Solution
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s)
    _accent_bar(s)
    _add_textbox(s, 0.7, 0.55, 10, 0.4, "OUR SOLUTION", size=14, bold=True, color=ACCENT)
    _add_textbox(
        s,
        0.7,
        1.0,
        11.5,
        1.0,
        "A controlled autonomous trading engine, not a black-box bot.",
        size=32,
        bold=True,
    )
    cols = [
        ("L0 Ingestion", "Market + external intelligence with retries, health checks, cache hydration."),
        ("L1 Finance Brain", "Scoring, spread modeling, confidence calibration, signal composition."),
        ("L2 Agent Panel", "Multi-agent thesis and intelligence verdicts (BUY / SKIP / WAIT)."),
        ("L3 / L4 Control", "Risk-checked execution + full telemetry and operator audit stream."),
    ]
    x_positions = [0.7, 3.5, 6.3, 9.1]
    for x, (title, body) in zip(x_positions, cols):
        _add_textbox(s, x, 2.4, 2.5, 0.5, title, size=16, bold=True, color=ACCENT2)
        _add_textbox(s, x, 2.95, 2.5, 3.8, body, size=13, color=MUTED)

    # 4 — Product
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s)
    _accent_bar(s)
    _add_textbox(s, 0.7, 0.55, 10, 0.4, "PRODUCT EXPERIENCE", size=14, bold=True, color=ACCENT)
    _add_textbox(
        s,
        0.7,
        1.0,
        11.5,
        0.9,
        "Bloomberg-style terminal for autonomous decision support.",
        size=32,
        bold=True,
    )
    products = [
        ("Arb Radar", "Ranked opportunities with net edge, confidence, settlement alignment."),
        ("Autopilot Console", "Live proposals, execution lifecycle, synchronized backend status."),
        ("Intelligence Layer", "Bright Data source checks, news risk, consensus overlays."),
        ("Operator Controls", "Paper-only enforcement, token-gated actions, deterministic fallbacks."),
    ]
    for i, (title, body) in enumerate(products):
        row, col = divmod(i, 2)
        _add_textbox(s, 0.7 + col * 6.0, 2.2 + row * 2.3, 5.5, 0.45, title, size=18, bold=True)
        _add_textbox(s, 0.7 + col * 6.0, 2.7 + row * 2.3, 5.5, 1.5, body, size=14, color=MUTED)

    # 5 — Traction
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s)
    _accent_bar(s)
    _add_textbox(s, 0.7, 0.55, 10, 0.4, "EXECUTION SNAPSHOT", size=14, bold=True, color=ACCENT)
    _add_textbox(
        s,
        0.7,
        1.0,
        11.5,
        0.9,
        "Running with real-time data pathways and strict safety defaults.",
        size=30,
        bold=True,
    )
    rows = [
        ("Arb opportunity generation", "Live", "100 rows", "/api/arb/opportunities"),
        ("Proposal generation", "Live", "18+ proposals", "/proposals"),
        ("Risk gates (paper mode)", "Enforced", "M01-first", "Risk check engine"),
        ("Frontend terminal", "Live", "Dashboard green", "Playwright verified"),
        ("Bright Data intelligence", "Active", "Verdict pipeline", "Agent + cron"),
    ]
    y = 2.15
    _add_textbox(s, 0.7, y, 3.2, 0.35, "Capability", size=13, bold=True, color=INK)
    _add_textbox(s, 3.9, y, 1.5, 0.35, "Status", size=13, bold=True, color=INK)
    _add_textbox(s, 5.4, y, 2.5, 0.35, "Signal", size=13, bold=True, color=INK)
    _add_textbox(s, 7.9, y, 4.0, 0.35, "Evidence", size=13, bold=True, color=INK)
    y += 0.45
    for cap, status, signal, ev in rows:
        _add_textbox(s, 0.7, y, 3.2, 0.35, cap, size=12, color=MUTED)
        _add_textbox(s, 3.9, y, 1.5, 0.35, status, size=12, color=OK)
        _add_textbox(s, 5.4, y, 2.5, 0.35, signal, size=12, color=MUTED)
        _add_textbox(s, 7.9, y, 4.0, 0.35, ev, size=11, color=MUTED)
        y += 0.55

    # 6 — Moat
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s)
    _accent_bar(s)
    _add_textbox(s, 0.7, 0.55, 10, 0.4, "WHY WE WIN", size=14, bold=True, color=ACCENT)
    _add_textbox(s, 0.7, 1.0, 11.5, 0.9, "Compound moat from architecture, controls, and learning loops.", size=32, bold=True)
    moat = [
        ("Safety-first by design", "Paper-only defaults and explicit gate ordering prevent accidental live exposure."),
        ("Cross-market intelligence", "Market microstructure plus verified external signals in one context."),
        ("Operator trust", "Transparent logs, deterministic fallbacks, visible rationale per action."),
        ("Data flywheel", "Historical calibration and model scoring improve quality over time."),
    ]
    for i, (title, body) in enumerate(moat):
        row, col = divmod(i, 2)
        _add_textbox(s, 0.7 + col * 6.0, 2.2 + row * 2.3, 5.5, 0.45, title, size=17, bold=True)
        _add_textbox(s, 0.7 + col * 6.0, 2.65 + row * 2.3, 5.5, 1.6, body, size=14, color=MUTED)

    # 7 — Roadmap
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s)
    _accent_bar(s)
    _add_textbox(s, 0.7, 0.55, 10, 0.4, "ROADMAP", size=14, bold=True, color=ACCENT)
    _add_textbox(
        s,
        0.7,
        1.0,
        11.5,
        0.9,
        "260-day plan from robust paper alpha to production-ready autonomy.",
        size=30,
        bold=True,
    )
    phases = [
        ("Phase 1", "Reliability + CI hardening, unified startup, regression prevention."),
        ("Phase 2", "Observability, operator workflows, decision explainability."),
        ("Phase 3", "Signal quality, settlement precision, stricter risk calibration."),
        ("Phase 4–5", "Scale architecture, auth, deployment hardening, staged autonomy."),
    ]
    for x, (title, body) in zip(x_positions, phases):
        _add_textbox(s, x, 2.3, 2.5, 0.45, title, size=16, bold=True, color=ACCENT2)
        _add_textbox(s, x, 2.8, 2.5, 3.5, body, size=13, color=MUTED)
    _add_textbox(
        s,
        0.7,
        6.2,
        11.5,
        0.6,
        "Themes: WC2026 brain · cron architecture · scheduled agents · trading upgrades · calibration seed",
        size=13,
        color=MUTED,
    )

    # 8 — Business model
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s)
    _accent_bar(s)
    _add_textbox(s, 0.7, 0.55, 10, 0.4, "BUSINESS MODEL", size=14, bold=True, color=ACCENT)
    _add_textbox(
        s,
        0.7,
        1.0,
        11.5,
        0.9,
        "From high-value operator tooling to fund-grade platform economics.",
        size=30,
        bold=True,
    )
    tiers = [
        ("Tier 1: Pro Terminal", "Subscription for autonomous radar, intelligence overlays, simulation analytics."),
        ("Tier 2: Team Ops", "Collaboration, policy controls, audit exports, strategy workspaces."),
        ("Tier 3: Institutional", "Dedicated deployments, custom integrations, risk policy frameworks."),
        ("Long-term", "Performance-linked products after live capital and compliance tracks."),
    ]
    for i, (title, body) in enumerate(tiers):
        row, col = divmod(i, 2)
        _add_textbox(s, 0.7 + col * 6.0, 2.2 + row * 2.3, 5.5, 0.45, title, size=17, bold=True)
        _add_textbox(s, 0.7 + col * 6.0, 2.65 + row * 2.3, 5.5, 1.6, body, size=14, color=MUTED)

    # 9 — Ask
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s)
    _accent_bar(s)
    _add_textbox(s, 0.7, 0.55, 10, 0.4, "THE ASK", size=14, bold=True, color=ACCENT)
    _add_textbox(
        s,
        0.7,
        1.0,
        11.5,
        1.0,
        "Partner with us to accelerate from strong engine to category-defining platform.",
        size=30,
        bold=True,
    )
    _add_textbox(s, 0.7, 2.5, 3.5, 0.8, "$X", size=40, bold=True)
    _add_textbox(s, 0.7, 3.25, 3.5, 0.5, "12-month build + GTM", size=14, color=MUTED)
    _add_textbox(s, 4.5, 2.5, 3.5, 0.8, "3 hires", size=40, bold=True)
    _add_textbox(s, 4.5, 3.25, 3.5, 0.5, "Infra · quant · product ops", size=14, color=MUTED)
    _add_textbox(s, 8.3, 2.5, 3.5, 0.8, "2 pilots", size=40, bold=True)
    _add_textbox(s, 8.3, 3.25, 3.5, 0.5, "Institutional design partners", size=14, color=MUTED)
    _add_textbox(
        s,
        0.7,
        4.2,
        11.5,
        1.2,
        "Use of funds: reliability and compliance infrastructure, alpha refinement, scalable hosted deployments.",
        size=18,
        color=MUTED,
    )

    # 10 — Close
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s)
    _accent_bar(s)
    _add_textbox(s, 0.7, 0.55, 10, 0.4, "CLOSING", size=14, bold=True, color=ACCENT)
    _add_textbox(s, 0.7, 1.2, 11.5, 1.2, "Autonomy with discipline wins this market.", size=44, bold=True)
    _add_textbox(
        s,
        0.7,
        2.6,
        11,
        1.0,
        "APEX Autopilot is building the operating system for intelligent prediction-market execution.",
        size=22,
        color=MUTED,
    )
    _add_textbox(
        s,
        0.7,
        4.0,
        10,
        1.0,
        "Team Autopilot — ready for product demo, technical diligence, and pilot planning.",
        size=20,
        color=INK,
    )

    prs.save(OUT)
    return OUT


if __name__ == "__main__":
    path = build()
    print(path)
